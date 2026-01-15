from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def create_presentation():
    prs = Presentation()

    # --- HELPER: Title Slide ---
    def add_title_slide(title, subtitle):
        slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(slide_layout)
        title_placeholder = slide.shapes.title
        subtitle_placeholder = slide.placeholders[1]
        title_placeholder.text = title
        subtitle_placeholder.text = subtitle

    # --- HELPER: Content Slide ---
    def add_content_slide(title, content_items):
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        title_placeholder = slide.shapes.title
        title_placeholder.text = title
        
        body_shape = slide.placeholders[1]
        tf = body_shape.text_frame
        
        for i, item in enumerate(content_items):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = item
            p.font.size = Pt(24)

    # --- SLIDE 1: Vision ---
    add_title_slide("ISAAC: The Semantic Bridge", "Building an AI-Ready Knowledge Network for BES")

    # --- SLIDE 2: The Challenge ---
    add_content_slide("The Current Landscape: Data Rich, Knowledge Poor", [
        "• We generate massive amounts of high-quality data at light sources.",
        "• But this data is trapped in silos (local drives, diverse formats).",
        "• AI Agents cannot 'reason' over raw pixels or proprietary binaries.",
        "• Result: We have a 'Data Lake', but not a 'Knowledge Graph'."
    ])

    # --- SLIDE 3: The Fear (Addressing User Concerns) ---
    add_content_slide("The Misconception", [
        "❌ \"Do I have to change how I do experiments?\"",
        "❌ \"Do I have to migrate my terabytes of raw data?\"",
        "❌ \"Will this slow down my beamtime?\"",
        " ",
        "👉 NO. The ISAAC Layer does not touch your raw data workflow."
    ])

    # --- SLIDE 4: The Reality (Architecture) ---
    add_content_slide("The Reality: ISAAC is an Overlay", [
        "• Tier 1: Raw Data stays where it is (Tape, Cluster, Local).",
        "• Tier 2: The ISAAC 'Harvester' runs post-experiment.",
        "• Tier 3: We generate a lightweight 'Knowledge Record' (KB).",
        " ",
        "• The Record is a 'Manifest' that points to the data, but describes the PHYSICS in a standard language."
    ])

    # --- SLIDE 5: AI-Readiness ---
    add_content_slide("What makes it 'AI-Ready'?", [
        "• Humans use ambiguity: 'sample=Cu' (Could be Copper or Control Unit).",
        "• AI needs precision: 'sample.material.formula = Cu'.",
        "• ISAAC provides a Controlled Vocabulary (The Ontology).",
        "• Result: An AI Agent can query 10,000 experiments across 5 facilities with 100% confidence."
    ])

    # --- SLIDE 6: The Vision ---
    # Placeholder for the visual
    slide_layout = prs.slide_layouts[5] # Blank
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "The Vision: Knowledge Islands"
    
    # Add Image if exists - Using the one we generated
    img_path = "/Users/dsokaras/.gemini/antigravity/brain/ede1e065-d118-473c-88f1-09cb74c02968/knowledge_islands_cluster_visualization_1767910553088.png"
    try:
        slide.shapes.add_picture(img_path, Inches(1), Inches(1.5), width=Inches(8))
    except:
        pass # Graceful fallback if image missing

    # --- SLIDE 7: Future State ---
    add_content_slide("The Future: From Files to Insights", [
        "• Today: You search for filenames.",
        "• Tomorrow: You ask the Agent...",
        "   \"Show me the trend in Faradaic Efficiency for all Copper catalysts tested in 2025.\"",
        "• The Agent traverses the 'Knowledge Islands' to build the answer.",
        "• This is the promise of the ISAAC Semantic Layer."
    ])

    # Save
    output_path = "ISAAC_Vision_Presentation.pptx"
    prs.save(output_path)
    print(f"Presentation saved to {output_path}")

if __name__ == "__main__":
    create_presentation()
